"""Export manager for generating reports."""

from pathlib import Path
from typing import Dict, Any, Optional
import pandas as pd
from datetime import datetime
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
import plotly.graph_objects as go
import plotly.io as pio


class ExportManager:
    """Manages report generation in various formats."""
    
    def __init__(self, output_dir: Optional[str] = None):
        """Initialize export manager.
        
        Args:
            output_dir: Output directory for reports
        """
        self.output_dir = Path(output_dir) if output_dir else Path("reports")
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate_pdf_report(
        self,
        df: pd.DataFrame,
        metadata: Dict[str, Any],
    ) -> str:
        """Generate PDF report.
        
        Args:
            df: Dataset to include
            metadata: Report metadata
            
        Returns:
            Path to generated PDF
        """
        filename = f"wastewater_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = self.output_dir / filename
        
        doc = SimpleDocTemplate(str(filepath), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f77b4'),
            spaceAfter=30,
        )
        story.append(Paragraph(metadata['title'], title_style))
        story.append(Paragraph(metadata['subtitle'], styles['Normal']))
        story.append(Paragraph(f"Generated: {metadata['date']}", styles['Normal']))
        story.append(Paragraph(f"Author: {metadata['author']}", styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        
        # Sections
        if 'Executive Summary' in metadata['sections']:
            story.append(Paragraph("Executive Summary", styles['Heading2']))
            story.append(Paragraph(
                f"This report analyzes {len(df)} records across {len(df.columns)} features.",
                styles['Normal']
            ))
            story.append(Spacer(1, 0.3*inch))
        
        if 'KPIs' in metadata['sections']:
            story.append(Paragraph("Key Performance Indicators", styles['Heading2']))
            numeric_cols = df.select_dtypes(include=['number']).columns[:5]
            
            kpi_data = [['Metric', 'Mean', 'Min', 'Max']]
            for col in numeric_cols:
                kpi_data.append([
                    col,
                    f"{df[col].mean():.2f}",
                    f"{df[col].min():.2f}",
                    f"{df[col].max():.2f}",
                ])
            
            table = Table(kpi_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            story.append(table)
            story.append(Spacer(1, 0.3*inch))
        
        if metadata['include_raw_data'] and len(df) < 100:
            story.append(Paragraph("Raw Data Sample", styles['Heading2']))
            # Add sample table (limit rows for PDF)
            sample_df = df.head(20)
            table_data = [list(sample_df.columns)] + sample_df.values.tolist()
            table = Table(table_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 8),
                ('FONTSIZE', (0, 1), (-1, -1), 7),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ]))
            story.append(table)
        
        doc.build(story)
        return str(filepath)
    
    def generate_pptx_report(
        self,
        df: pd.DataFrame,
        metadata: Dict[str, Any],
    ) -> str:
        """Generate PowerPoint report.
        
        Args:
            df: Dataset to include
            metadata: Report metadata
            
        Returns:
            Path to generated PPTX
        """
        filename = f"wastewater_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = self.output_dir / filename
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = metadata['title']
        subtitle.text = f"{metadata['subtitle']}\n{metadata['date']}\n{metadata['author']}"
        
        # KPI slide
        if 'KPIs' in metadata['sections']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Key Performance Indicators"
            
            # Add table with KPIs
            numeric_cols = df.select_dtypes(include=['number']).columns[:5]
            rows = min(6, len(numeric_cols) + 1)
            cols = 4
            
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            table.columns[0].width = Inches(2)
            
            # Header row
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Mean"
            table.cell(0, 2).text = "Min"
            table.cell(0, 3).text = "Max"
            
            # Data rows
            for i, col in enumerate(numeric_cols[:5], 1):
                if i < rows:
                    table.cell(i, 0).text = col[:20]
                    table.cell(i, 1).text = f"{df[col].mean():.2f}"
                    table.cell(i, 2).text = f"{df[col].min():.2f}"
                    table.cell(i, 3).text = f"{df[col].max():.2f}"
        
        prs.save(str(filepath))
        return str(filepath)
    
    def generate_html_report(
        self,
        df: pd.DataFrame,
        metadata: Dict[str, Any],
    ) -> str:
        """Generate HTML report.
        
        Args:
            df: Dataset to include
            metadata: Report metadata
            
        Returns:
            Path to generated HTML
        """
        filename = f"wastewater_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        filepath = self.output_dir / filename
        
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <title>{metadata['title']}</title>
            <style>
                body {{ font-family: Inter, sans-serif; margin: 40px; background: #f8f9fa; }}
                .header {{ background: linear-gradient(135deg, #1f77b4 0%, #2a5f8f 100%); color: white; padding: 30px; border-radius: 8px; }}
                .section {{ background: white; padding: 20px; margin: 20px 0; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }}
                table {{ width: 100%; border-collapse: collapse; }}
                th {{ background: #1f77b4; color: white; padding: 12px; text-align: left; }}
                td {{ padding: 10px; border-bottom: 1px solid #ddd; }}
                tr:hover {{ background: #f5f5f5; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>{metadata['title']}</h1>
                <p>{metadata['subtitle']}</p>
                <p>Generated: {metadata['date']} | Author: {metadata['author']}</p>
            </div>
            
            <div class="section">
                <h2>Dataset Overview</h2>
                <p>Total Records: {len(df):,}</p>
                <p>Features: {len(df.columns)}</p>
            </div>
        """
        
        if 'KPIs' in metadata['sections']:
            numeric_cols = df.select_dtypes(include=['number']).columns[:10]
            html_content += """
            <div class="section">
                <h2>Key Performance Indicators</h2>
                <table>
                    <thead>
                        <tr>
                            <th>Metric</th>
                            <th>Mean</th>
                            <th>Min</th>
                            <th>Max</th>
                            <th>Std Dev</th>
                        </tr>
                    </thead>
                    <tbody>
            """
            for col in numeric_cols:
                html_content += f"""
                        <tr>
                            <td>{col}</td>
                            <td>{df[col].mean():.2f}</td>
                            <td>{df[col].min():.2f}</td>
                            <td>{df[col].max():.2f}</td>
                            <td>{df[col].std():.2f}</td>
                        </tr>
                """
            html_content += """
                    </tbody>
                </table>
            </div>
            """
        
        html_content += """
        </body>
        </html>
        """
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return str(filepath)



